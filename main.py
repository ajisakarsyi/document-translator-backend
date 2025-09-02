import os
import io
import re
import json
import tempfile
from enum import Enum
from typing import Dict, Any, Optional

from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Depends
from fastapi.responses import StreamingResponse, JSONResponse, PlainTextResponse

from pptx import Presentation
from docx import Document
import fitz  # PyMuPDF
import requests
from dotenv import load_dotenv
from pdf2docx import Converter

# Import authentication system
from auth import require_auth, IntrospectionResult

load_dotenv()

app = FastAPI(title="Translation API", version="1.0.0")

# =========================
# Helpers
# =========================
def require_role(principal: IntrospectionResult, required_role: str):
    if not principal.role or principal.role.lower() != required_role.lower():
        raise HTTPException(
            status_code=403,
            detail=f"Requires {required_role} role"
        )

# In-memory storage (for demo purposes)
pending_requests: Dict[str, dict] = {}
approved_docs: Dict[str, dict] = {}

# =========================
# Health / test endpoints
# =========================
@app.get("/health", response_class=PlainTextResponse)
def health():
    return "ok"


@app.get("/health/auth")
async def health_auth(principal: IntrospectionResult = Depends(require_auth)):
    return {
        "status": "ok",
        "sub": principal.sub,
        "aud": principal.aud,
        "scope": principal.scope,
        "role": principal.role,
    }


@app.get("/whoami")
async def whoami(principal: IntrospectionResult = Depends(require_auth)):
    return JSONResponse({
        "sub": principal.sub,
        "aud": principal.aud,
        "scope": principal.scope,
        "client_id": principal.client_id,
        "role": principal.role,
    })

# =========================
# Role-based document flows
# =========================
@app.post("/user/upload")
async def user_upload(
    file: UploadFile = File(...),
    principal: IntrospectionResult = Depends(require_auth)
):
    require_role(principal, "user")

    request_id = f"req-{len(pending_requests)+1}"
    pending_requests[request_id] = {
        "filename": file.filename,
        "uploaded_by": principal.sub,
        "status": "pending"
    }
    return {"request_id": request_id, "status": "pending"}


@app.post("/admin/upload")
async def admin_upload(
    file: UploadFile = File(...),
    principal: IntrospectionResult = Depends(require_auth)
):
    require_role(principal, "admin")

    doc_id = f"doc-{len(approved_docs)+1}"
    approved_docs[doc_id] = {
        "filename": file.filename,
        "uploaded_by": principal.sub,
        "status": "approved"
    }
    return {"doc_id": doc_id, "status": "approved"}


@app.post("/admin/approve/{request_id}")
async def approve_request(
    request_id: str,
    principal: IntrospectionResult = Depends(require_auth)
):
    require_role(principal, "admin")

    if request_id not in pending_requests:
        raise HTTPException(404, "Request not found")

    pending_requests[request_id]["status"] = "approved"
    approved_docs[request_id] = pending_requests.pop(request_id)

    return {"request_id": request_id, "status": "approved"}

# =========================
# Translation logic
# =========================
OLLAMA_BASE_URL = os.getenv("OLLAMA_BASE_URL", "http://localhost:11434")
OLLAMA_MODEL_NAME = os.getenv("OLLAMA_MODEL_NAME", "mistral-nemo:12b")


class Language(str, Enum):
    english = "English"
    japanese = "Japanese"
    indonesian = "Indonesian"
    french = "French"
    spanish = "Spanish"
    german = "German"


def translate_text_ollama(text: str, source_language: str, target_language: str) -> str:
    prompt = f"""
Translate the following text from {source_language} to {target_language}.
Keep placeholders, brand names, and technical terms unchanged.
Preserve line breaks and formatting.
Do not return your thinking process, internal notes, or explanations.
ONLY the translated text directly.

Text:
{text}

Translated:
"""
    try:
        response = requests.post(
            f"{OLLAMA_BASE_URL}/api/generate",
            json={"model": OLLAMA_MODEL_NAME, "prompt": prompt, "stream": False},
            timeout=120,
        )
        response.raise_for_status()
        raw_translation = response.json().get("response", "").strip()

        cleaned = re.sub(r"<think>.*?</think>", "", raw_translation, flags=re.DOTALL)
        cleaned = re.sub(r"\*\*Note[s]?:\*\*.*", "", cleaned, flags=re.DOTALL)
        return cleaned.strip()
    except Exception as e:
        print(f"Ollama Error: {e}")
        return f"[Translation Failed] {text}"


def extract_text_from_pdf(file_bytes: bytes):
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    blocks, layout = [], []

    for page in doc:
        texts = page.get_text("blocks")
        page_blocks = []

        if not texts:
            plain_text = page.get_text("text").strip()
            if plain_text:
                blocks.append(plain_text)
                layout.append({"images": [], "texts": [(plain_text, (0, 0, 500, 500))]})
                continue

        for b in texts:
            if len(b) < 7:
                continue
            x0, y0, x1, y1, text, _, block_type = b
            if block_type == 0 and text.strip():
                norm_text = text.strip()
                blocks.append(norm_text)
                page_blocks.append((norm_text, (x0, y0, x1, y1)))
        layout.append({"images": [], "texts": page_blocks})

    doc.close()

    if not blocks:
        raise ValueError("No text could be extracted from PDF. It may be scanned or image-based.")

    return "\n\n".join(blocks), layout


def extract_text_from_ppt(file_stream: io.BytesIO):
    prs = Presentation(file_stream)
    blocks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                blocks.append(shape.text.strip())
    return "\n\n".join(blocks)


def convert_pdf_to_docx(pdf_bytes: bytes) -> io.BytesIO:
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
        temp_pdf.write(pdf_bytes)
        temp_pdf_path = temp_pdf.name

    temp_docx_path = temp_pdf_path.replace(".pdf", ".docx")

    try:
        cv = Converter(temp_pdf_path)
        cv.convert(temp_docx_path)
        cv.close()

        docx_buffer = io.BytesIO()
        with open(temp_docx_path, "rb") as f:
            docx_buffer.write(f.read())
        docx_buffer.seek(0)

        return docx_buffer

    finally:
        try:
            os.remove(temp_pdf_path)
        except Exception as e:
            print(f"[WARN] Could not delete temp PDF: {e}")
        try:
            os.remove(temp_docx_path)
        except Exception as e:
            print(f"[WARN] Could not delete temp DOCX: {e}")


def translate_docx(file_stream: io.BytesIO, source_language: str, target_language: str):
    doc = Document(file_stream)
    for para in doc.paragraphs:
        if para.text.strip():
            translated = translate_text_ollama(para.text.strip(), source_language, target_language)
            para.text = translated
    out_stream = io.BytesIO()
    doc.save(out_stream)
    out_stream.seek(0)
    return out_stream


@app.post("/translate-document/", response_class=StreamingResponse)
async def translate_document(
    file: UploadFile = File(...),
    source_language: Language = Form(...),
    target_language: Language = Form(...),
    principal: IntrospectionResult = Depends(require_auth)  # <-- Auth enforced
):
    filename = file.filename
    ext = os.path.splitext(filename)[1].lower()
    content = await file.read()
    source = source_language.value
    target = target_language.value

    if source == target:
        raise HTTPException(400, "Source and target languages must differ.")

    if ext == ".pdf":
        try:
            docx_buffer = convert_pdf_to_docx(content)
            translated_stream = translate_docx(docx_buffer, source, target)
            return StreamingResponse(
                translated_stream,
                headers={"Content-Disposition": f"attachment; filename=translated_{filename.replace('.pdf', '.docx')}"},
                media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            )
        except Exception as e:
            raise HTTPException(500, f"PDF translation failed: {e}")

    elif ext == ".pptx":
        text = extract_text_from_ppt(io.BytesIO(content))
        translated = translate_text_ollama(text, source, target)

        prs = Presentation(io.BytesIO(content))
        lines = translated.split("\n\n")
        idx = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if idx < len(lines):
                        shape.text = lines[idx]
                        idx += 1
        out_ppt = io.BytesIO()
        prs.save(out_ppt)
        out_ppt.seek(0)
        return StreamingResponse(
            out_ppt,
            headers={"Content-Disposition": f"attachment; filename=translated_{filename}"},
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    elif ext == ".docx":
        translated_stream = translate_docx(io.BytesIO(content), source, target)
        return StreamingResponse(
            translated_stream,
            headers={"Content-Disposition": f"attachment; filename=translated_{filename}"},
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )

    else:
        raise HTTPException(400, "Only PDF, PPTX, and DOCX formats are supported.")
